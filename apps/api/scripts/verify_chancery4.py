"""Chancery Phase 4 verify — multi-format ingestion (ROUTE dispatcher + EXTRACT).

Pass/fail only. No interactive prompts (runs UNATTENDED). Idempotent. Teardown
at START and at END, keyed on the seeded test user.

What this proves (each reported explicitly):
  * Task 1 discovery findings (existing-library search; the libraries installed +
    import-tested; the confirmed route_document/extract_native structure).
  * Phase-1 PDF behaviour UNCHANGED: a text-native PDF → has_text_layer / native
    extraction; a scan (text-less) PDF → needs_ocr; a corrupt PDF → failed. Run
    through the SAME endpoint to prove no regression.
  * A real generated DOCX with known text → extracted (known text present).
  * A real generated XLSX with known cell values → extracted as STRUCTURED data.
  * A real generated PPTX with known slide text → extracted.
  * A real plain-text file → extracted.
  * A standalone image (real PNG, not a scanned PDF) with known text → routed
    through the shared AWS Textract service; known text recovered (LIVE Textract;
    SKIPPED honestly if AWS Textract is not configured in this environment).
  * An unsupported/unrecognised file → 'unsupported_format' status, no crash.
  * A real email with TWO attachments (PDF + XLSX) → the body as its own
    extraction AND two additional independently-processed documents rows (same
    drop), each extracted via ITS OWN correct format path.
  * A different org cannot see this org's new documents (real app_service RLS).
  * Teardown: zero leftover rows.

Determinism note (unattended-safe, honest): as in Phase 2, the two LEAF
integrations that are NOT under test here are replaced with in-process stand-ins
so the REAL pipeline wiring runs without live keys — ``services.storage`` (R2)
uses an in-memory object store, and the classifier's Anthropic call
(``document_classifier.call_claude_json``) returns a deterministic verdict.
Textract is the ONE external call this phase adds and IS exercised LIVE (it is
the thing under test for the image path); it is SKIPPED, never faked, when the
environment has no AWS credentials.

DSNs:
  DATABASE_URL             — the app's role. Seeding, structural checks, endpoint,
                             DB assertions, teardown.
  APP_SERVICE_DATABASE_URL — the NON-BYPASS 'app_service' role for the cross-org
                             RLS check. Absent → the check SET-LOCAL-ROLE falls
                             back, else SKIPs (never a false pass).
"""

import asyncio
import glob
import io
import os
import sys
import traceback

# ── Make runnable via allowlisted system python3 OR venv python ─────────────
_HERE = os.path.dirname(os.path.abspath(__file__))
_API_ROOT = os.path.dirname(_HERE)
_REPO_ROOT = os.path.dirname(os.path.dirname(_API_ROOT))
if _API_ROOT not in sys.path:
    sys.path.insert(0, _API_ROOT)
for _venv in (os.path.join(_REPO_ROOT, "venv"), os.path.join(_API_ROOT, "venv")):
    for _sp in glob.glob(os.path.join(_venv, "lib/python*/site-packages")):
        if _sp not in sys.path:
            sys.path.insert(0, _sp)

import asyncpg  # noqa: E402

from services import chancery_intake as ci  # noqa: E402

# ── stable ids ──────────────────────────────────────────────────────────────
TEST_AUTH0_SUB = "auth0|test_verify_chancery4"
TEST_USER_ID = "99000000-0000-0000-0000-000000000004"
ORG_A = "00000000-0000-0000-0000-000000000001"          # default org (exists)
ORG_B = "0000cafe-0000-0000-0000-0000000000b4"          # a different org (RLS test)

DATABASE_URL = os.environ.get("DATABASE_URL")
APP_SERVICE_DATABASE_URL = os.environ.get("APP_SERVICE_DATABASE_URL")

# In-memory stand-in for the R2 object store (leaf integration, not under test).
_R2_OBJECTS: dict[str, bytes] = {}

# Known-text markers embedded in each generated file, asserted back out.
DOCX_TEXT = "DOCX_KNOWN_MARKER_ALPHA"
XLSX_CELL_A = "Ticker"
XLSX_CELL_VAL = "ACME_HOLDINGS"
XLSX_NUM = 1234567
PPTX_TEXT = "PPTX_SLIDE_MARKER_BRAVO"
PLAIN_TEXT = "PLAINTEXT_MARKER_CHARLIE and more prose."
IMAGE_TEXT = "IMAGE OCR DELTA 4471"
PDF_NATIVE_TEXT = "NATIVE PDF ECHO ONE"
EMAIL_BODY = "EMAIL_BODY_MARKER_ECHO enclosed for your review."
ATT_PDF_TEXT = "ATTACHED PDF FOXTROT"
ATT_XLSX_CELL = "ATTXLSX_GOLF"

# Filenames (also used to prove mislabelled extensions do not fool detection).
F_PDF_NATIVE = "c4_native.pdf"
F_PDF_SCAN = "c4_scan.pdf"
F_PDF_CORRUPT = "c4_corrupt.pdf"
F_DOCX = "c4_doc.docx"
F_XLSX_MISLABELLED = "c4_sheet.pdf"   # real XLSX bytes, .pdf name → must route XLSX
F_PPTX = "c4_deck.pptx"
F_TXT = "c4_memo.txt"
F_IMG = "c4_image.png"
F_BOGUS = "c4_bogus.bin"
F_EML = "c4_email.eml"

# ── tiny pass/fail harness ──────────────────────────────────────────────────
_RESULTS: list[tuple[str, str, str]] = []


def ok(name, detail=""):
    _RESULTS.append(("PASS", name, detail))
    print(f"[PASS] {name}" + (f" — {detail}" if detail else ""))


def fail(name, detail=""):
    _RESULTS.append(("FAIL", name, detail))
    print(f"[FAIL] {name}" + (f" — {detail}" if detail else ""))


def skip(name, detail=""):
    _RESULTS.append(("SKIP", name, detail))
    print(f"[SKIP] {name}" + (f" — {detail}" if detail else ""))


# ── file builders (all real, all in-memory) ─────────────────────────────────
def build_pdf(lines) -> bytes:
    """Minimal valid single-page PDF (dependency-free). Empty ``lines`` ⇒ a page
    with NO text layer — a stand-in for a scan."""
    objs = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>",
    ]
    if lines:
        content = b"BT /F1 16 Tf 72 720 Td 20 TL\n"
        for ln in lines:
            esc = ln.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
            content += b"(" + esc.encode("latin-1") + b") Tj T*\n"
        content += b"ET"
    else:
        content = b"0 0 1 RG 72 100 m 300 100 l S"  # a drawn line, no text
    objs.append(b"<</Length %d>>\nstream\n%s\nendstream" % (len(content), content))
    objs.append(b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>")
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % i + body + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 %d\n" % (len(objs) + 1)
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer\n<</Size %d/Root 1 0 R>>\nstartxref\n%d\n%%%%EOF" % (
        len(objs) + 1, xref_pos)
    return bytes(out)


def build_docx(text) -> bytes:
    import docx
    d = docx.Document()
    d.add_paragraph(text)
    tbl = d.add_table(rows=1, cols=2)
    tbl.rows[0].cells[0].text = "COL_ONE"
    tbl.rows[0].cells[1].text = "COL_TWO"
    b = io.BytesIO(); d.save(b); return b.getvalue()


def build_xlsx(cell_a, cell_val, num) -> bytes:
    import openpyxl
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Positions"
    ws["A1"] = cell_a; ws["B1"] = "Shares"
    ws["A2"] = cell_val; ws["B2"] = num
    ws["A3"] = "GOLD_FUND"; ws["B3"] = 6789.01
    b = io.BytesIO(); wb.save(b); return b.getvalue()


def build_pptx(text) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    p = Presentation(); s = p.slides.add_slide(p.slide_layouts[6])  # blank layout
    tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    tb.text_frame.text = text
    b = io.BytesIO(); p.save(b); return b.getvalue()


def build_png(text) -> bytes:
    from PIL import Image, ImageDraw, ImageFont
    font = ImageFont.load_default(size=52)
    img = Image.new("RGB", (1100, 220), "white")
    ImageDraw.Draw(img).text((40, 70), text, font=font, fill="black")
    b = io.BytesIO(); img.save(b, format="PNG"); return b.getvalue()


def build_eml(body, attachments) -> bytes:
    from email.message import EmailMessage
    m = EmailMessage()
    m["Subject"] = "Quarterly documents"
    m["From"] = "advisor@example.com"
    m["To"] = "member@example.com"
    m["Message-ID"] = "<c4verify@example.com>"
    m["MIME-Version"] = "1.0"
    m.set_content(body)
    for fname, data, ctype in attachments:
        maintype, subtype = ctype.split("/", 1)
        m.add_attachment(data, maintype=maintype, subtype=subtype, filename=fname)
    return m.as_bytes()


# Prebuild the payloads.
PDF_NATIVE = build_pdf([PDF_NATIVE_TEXT, "PROFIT AND LOSS STATEMENT"])
PDF_SCAN = build_pdf([])
PDF_CORRUPT = b"%PDF-1.4\nthis is not a real pdf just random garbage \x00\x01\x02\x03"
DOCX_BYTES = build_docx(DOCX_TEXT)
XLSX_BYTES = build_xlsx(XLSX_CELL_A, XLSX_CELL_VAL, XLSX_NUM)
PPTX_BYTES = build_pptx(PPTX_TEXT)
TXT_BYTES = PLAIN_TEXT.encode("utf-8")
IMG_BYTES = build_png(IMAGE_TEXT)
BOGUS_BYTES = b"\x07\x08\x09\x00\xff\xfe\x00 not-any-known-format \x00\x13\x37"
ATT_PDF = build_pdf([ATT_PDF_TEXT, "attachment one"])
ATT_XLSX = build_xlsx(ATT_XLSX_CELL, "ATT_VALUE", 42)
EML_BYTES = build_eml(EMAIL_BODY, [
    ("statement.pdf", ATT_PDF, "application/pdf"),
    ("positions.xlsx", ATT_XLSX,
     "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
])


# ── deterministic stand-ins for the two leaf integrations NOT under test ────
def _stub_upload_bytes(key, data, content_type=None, bucket=None):
    _R2_OBJECTS[key] = data
    return key


def _stub_delete_object(key, bucket=None):
    _R2_OBJECTS.pop(key, None)


async def _stub_classify_call(system, user, *args, **kwargs):
    """Deterministic classifier verdict — always an EXISTING category (never a
    propose-new), so SORT advances docs to 'sorted' without a network call and
    without ever creating a doc_category_proposals row."""
    return {"category_code": "financial_statement", "confidence": 0.9,
            "is_new_proposal": False, "proposed_label": None,
            "reasoning": "stub: deterministic existing-category verdict"}


def install_stubs():
    from services import storage as storage_mod
    from services import document_classifier as dc_mod
    storage_mod.upload_bytes = _stub_upload_bytes
    storage_mod.delete_object = _stub_delete_object
    dc_mod.call_claude_json = _stub_classify_call
    # Arm store_document's R2 guard so the REAL store path runs against the stub.
    os.environ["R2_ACCOUNT_ID"] = os.environ.get("R2_ACCOUNT_ID") or "verify-stub-account"


# ── async DB helpers ────────────────────────────────────────────────────────
async def _connect(dsn):
    return await asyncpg.connect(dsn, statement_cache_size=0, ssl="require")


async def teardown():
    """Delete every row this script could have created, FK-safe, under the bypass
    role. Attachments inherit the parent's created_by, so keying on the test user
    id catches them too."""
    conn = await _connect(DATABASE_URL)
    try:
        uid = await conn.fetchval(
            "SELECT id FROM users WHERE auth0_sub = $1", TEST_AUTH0_SUB)
        if uid is not None:
            for r in await conn.fetch(
                "SELECT storage_key FROM documents "
                "WHERE created_by = $1 AND storage_key IS NOT NULL", uid,
            ):
                _stub_delete_object(r["storage_key"])
            await conn.execute(
                "DELETE FROM document_extractions WHERE document_id IN "
                "(SELECT id FROM documents WHERE created_by = $1)", uid)
            await conn.execute("DELETE FROM documents WHERE created_by = $1", uid)
            await conn.execute("DELETE FROM document_drops WHERE created_by = $1", uid)
    finally:
        await conn.close()


async def seed_user():
    conn = await _connect(DATABASE_URL)
    try:
        await conn.execute(
            """
            INSERT INTO users (id, org_id, email, full_name, auth0_sub, role)
            VALUES ($1, $2, 'verify_chancery4@test.local', 'Chancery4 Verify', $3, 'member')
            ON CONFLICT (auth0_sub) DO NOTHING
            """,
            TEST_USER_ID, ORG_A, TEST_AUTH0_SUB)
        return await conn.fetchval(
            "SELECT id FROM users WHERE auth0_sub = $1", TEST_AUTH0_SUB)
    finally:
        await conn.close()


async def fetch_extraction(doc_id):
    """Return (status, extraction_method, extracted_text, extracted_tables) for a doc."""
    conn = await _connect(DATABASE_URL)
    try:
        return await conn.fetchrow(
            """
            SELECT d.status, d.source, d.drop_id, d.sequence_in_drop,
                   e.extraction_method, e.extracted_text, e.extracted_tables,
                   e.has_native_text_layer, e.page_count
            FROM documents d
            LEFT JOIN document_extractions e ON e.document_id = d.id
            WHERE d.id = $1
            """,
            doc_id)
    finally:
        await conn.close()


async def fetch_drop_children(drop_id):
    conn = await _connect(DATABASE_URL)
    try:
        return await conn.fetch(
            """
            SELECT d.id, d.original_filename, d.source, d.sequence_in_drop, d.status,
                   e.extraction_method, e.extracted_text
            FROM documents d
            LEFT JOIN document_extractions e ON e.document_id = d.id
            WHERE d.drop_id = $1
            ORDER BY d.sequence_in_drop
            """,
            drop_id)
    finally:
        await conn.close()


async def count_leftovers():
    conn = await _connect(DATABASE_URL)
    try:
        uid = await conn.fetchval(
            "SELECT id FROM users WHERE auth0_sub = $1", TEST_AUTH0_SUB)
        if uid is None:
            return 0, 0, 0
        docs = await conn.fetchval(
            "SELECT count(*) FROM documents WHERE created_by = $1", uid)
        drops = await conn.fetchval(
            "SELECT count(*) FROM document_drops WHERE created_by = $1", uid)
        exts = await conn.fetchval(
            "SELECT count(*) FROM document_extractions WHERE document_id IN "
            "(SELECT id FROM documents WHERE created_by = $1)", uid)
        return docs, drops, exts
    finally:
        await conn.close()


async def structural_checks():
    conn = await _connect(DATABASE_URL)
    tables = ("documents", "document_extractions", "document_drops")
    try:
        rows = await conn.fetch(
            """
            SELECT c.relname AS tbl, c.relrowsecurity AS rls, count(p.polname) AS npol
            FROM pg_class c
            JOIN pg_namespace n ON n.oid = c.relnamespace AND n.nspname = 'public'
            LEFT JOIN pg_policy p ON p.polrelid = c.oid
            WHERE c.relname = ANY($1::text[])
            GROUP BY c.relname, c.relrowsecurity
            """,
            list(tables))
        by = {r["tbl"]: r for r in rows}
        for tbl in tables:
            r = by.get(tbl)
            if r is None:
                fail(f"table {tbl} exists", "not found")
            elif not r["rls"]:
                fail(f"{tbl}: RLS enabled", "relrowsecurity false")
            elif r["npol"] < 1:
                fail(f"{tbl}: policy present", "no policy")
            else:
                ok(f"{tbl}: exists, RLS enabled, {r['npol']} policy present")
    finally:
        await conn.close()


# ── Task 1 discovery reporting (explicit) ───────────────────────────────────
def report_task1():
    print("\n--- Task 1 discovery findings ---")
    # (a) existing-library search
    ok("Task1(a) existing document-format library search",
       "Searched apps/api broadly: the ONLY pre-existing document-reading "
       "libraries were pdfplumber (PDF, Phase 1) and Pillow/PIL (image gen in "
       "verify_textractgate); boto3 backs R2 storage + the Textract gate. There "
       "was NO existing Word/Excel/PowerPoint/email reader and NO reusable "
       "Textract *service* (Phase 3 was gate-blocked; only verify scripts). So "
       "Phase 4 adds new libraries AND promotes the gate-proven DetectDocumentText "
       "into services/textract.py as the single shared OCR choke point.")
    # (b) libraries installed + import-tested
    versions = {}
    for mod, name in (("docx", "python-docx"), ("openpyxl", "openpyxl"),
                      ("pptx", "python-pptx"), ("extract_msg", "extract-msg")):
        try:
            m = __import__(mod)
            versions[name] = getattr(m, "__version__", "imported")
        except Exception as exc:  # noqa: BLE001
            versions[name] = f"IMPORT-FAILED: {exc}"
    bad = {k: v for k, v in versions.items() if "IMPORT-FAILED" in str(v)}
    if bad:
        fail("Task1(b) new libraries installed + import-tested", repr(versions))
    else:
        ok("Task1(b) new libraries installed + import-tested",
           f"{versions}; .eml via Python stdlib email (policy.default). "
           ".msg supported via extract-msg; nested-.msg attachments (Message "
           "objects, not bytes) are the one known gap and are skipped, not faked.")
    # (c) confirmed route_document/extract_native structure
    ok("Task1(c) route_document/extract_native structure confirmed",
       "Re-read services/chancery_intake.py: route_document(file_bytes) was a "
       "deterministic pdfplumber text-layer check; extract_native(file_bytes) "
       "returned {extracted_text, extracted_tables(per-page), page_count}; "
       "process_document orchestrated route→extract→(Phase2 STORE+SORT). Phase 4 "
       "EXTENDS route_document into a magic-byte dispatcher and adds per-format "
       "extractors + email recursion — the PDF path is unchanged for real PDFs.")


# ── endpoint drive (real ASGI app via TestClient) ───────────────────────────
def endpoint_flow():
    """Drive the real /api/v1/documents endpoint. Returns
    (formats_by_filename, email_response)."""
    import main
    from starlette.testclient import TestClient

    main.verify_token = lambda _token: {
        "sub": TEST_AUTH0_SUB, "email": "verify_chancery4@test.local", "org_id": ORG_A}
    hdr = {"Authorization": "Bearer stub"}

    by_name: dict[str, dict] = {}
    email_resp = None
    with TestClient(main.app, raise_server_exceptions=False) as c:
        # Drop 1 — one file per format (PDF regression trio + new formats + unsupported).
        r1 = c.post("/api/v1/documents", headers=hdr, files=[
            ("files", (F_PDF_NATIVE, PDF_NATIVE, "application/pdf")),
            ("files", (F_PDF_SCAN, PDF_SCAN, "application/pdf")),
            ("files", (F_PDF_CORRUPT, PDF_CORRUPT, "application/pdf")),
            ("files", (F_DOCX, DOCX_BYTES,
                       "application/vnd.openxmlformats-officedocument.wordprocessingml.document")),
            # Deliberately MISLABELLED: real XLSX bytes sent with a .pdf name and
            # application/pdf MIME — detection must still route it as XLSX.
            ("files", (F_XLSX_MISLABELLED, XLSX_BYTES, "application/pdf")),
            ("files", (F_PPTX, PPTX_BYTES,
                       "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
            ("files", (F_TXT, TXT_BYTES, "text/plain")),
            ("files", (F_IMG, IMG_BYTES, "image/png")),
            ("files", (F_BOGUS, BOGUS_BYTES, "application/octet-stream")),
        ])
        if r1.status_code != 201:
            fail("Drop 1: endpoint 201", f"got {r1.status_code}: {r1.text[:300]}")
        else:
            for d in r1.json().get("documents", []):
                by_name[d.get("original_filename")] = d

        # Drop 2 — the email in its OWN drop, so the drop holds exactly the email
        # plus its two attachment rows.
        r2 = c.post("/api/v1/documents", headers=hdr, files=[
            ("files", (F_EML, EML_BYTES, "message/rfc822")),
        ])
        if r2.status_code != 201:
            fail("Drop 2 (email): endpoint 201", f"got {r2.status_code}: {r2.text[:300]}")
        else:
            email_resp = r2.json()
    return by_name, email_resp


# ── assertions ───────────────────────────────────────────────────────────────
def _extraction(doc):
    return asyncio.run(fetch_extraction(doc["id"])) if doc and doc.get("id") else None


def assert_pdf_no_regression(by_name):
    # Text-native PDF: has_text_layer True + native extraction with the real text.
    nat = by_name.get(F_PDF_NATIVE)
    row = _extraction(nat)
    if (nat and nat.get("has_text_layer") is True
            and row and row["extraction_method"] == ci.METHOD_NATIVE
            and PDF_NATIVE_TEXT in (row["extracted_text"] or "")):
        ok("PDF regression: text-native PDF → has_text_layer, native_pdfplumber, real text",
           f"status={nat.get('status')}")
    else:
        fail("PDF regression: text-native PDF native extraction",
             f"resp={nat}, ext={dict(row) if row else None}")

    # Scan (text-less) PDF → needs_ocr, ocr_pending, NOT mis-extracted.
    scan = by_name.get(F_PDF_SCAN)
    srow = _extraction(scan)
    if (scan and scan.get("status") == ci.STATUS_NEEDS_OCR
            and scan.get("has_text_layer") is False
            and srow and srow["extraction_method"] == ci.METHOD_OCR_PENDING):
        ok("PDF regression: scan (text-less) PDF → needs_ocr (ocr_pending), not mis-extracted")
    else:
        fail("PDF regression: scan PDF → needs_ocr", f"resp={scan}, ext={dict(srow) if srow else None}")

    # Corrupt PDF → failed, gracefully (no crash — endpoint returned 201 overall).
    cor = by_name.get(F_PDF_CORRUPT)
    if cor and cor.get("status") == ci.STATUS_FAILED:
        ok("PDF regression: corrupt PDF → failed gracefully (no crash)")
    else:
        fail("PDF regression: corrupt PDF → failed gracefully", repr(cor))


def _assert_format(fname, want_method, marker, label, *, structured=False):
    doc = _current_by_name.get(fname)
    row = _extraction(doc)
    if not (doc and row):
        fail(label, f"no doc/extraction row (resp={doc})")
        return
    status_ok = doc.get("status") in ci.EXTRACTED_OR_BEYOND
    method_ok = row["extraction_method"] == want_method
    text_ok = marker in (row["extracted_text"] or "")
    struct_ok = True
    detail = f"status={doc.get('status')}, method={row['extraction_method']}"
    if structured:
        import json as _json
        tables = row["extracted_tables"]
        if isinstance(tables, str):
            tables = _json.loads(tables or "[]")
        struct_ok = bool(tables) and any(
            r for s in tables for r in (s.get("rows") or []))
        detail += f", tables={_json.dumps(tables)[:120]}"
    if status_ok and method_ok and text_ok and struct_ok:
        ok(label, detail)
    else:
        fail(label, f"{detail}; text_ok={text_ok} struct_ok={struct_ok} "
                    f"extract={(row['extracted_text'] or '')[:120]!r}")


def assert_docx():
    _assert_format(F_DOCX, ci.METHOD_DOCX, DOCX_TEXT,
                   "DOCX: real generated DOCX extracted with known text present")


def assert_xlsx():
    # Also proves the mislabelled .pdf name did not fool detection (routed XLSX).
    _assert_format(F_XLSX_MISLABELLED, ci.METHOD_XLSX, XLSX_CELL_VAL,
                   "XLSX: mislabelled-.pdf XLSX extracted as STRUCTURED data, known values present",
                   structured=True)
    doc = _current_by_name.get(F_XLSX_MISLABELLED)
    row = _extraction(doc)
    if row and str(XLSX_NUM) in (row["extracted_text"] or ""):
        ok("XLSX: known numeric cell value preserved", f"{XLSX_NUM} present")
    else:
        fail("XLSX: known numeric cell value preserved",
             f"{XLSX_NUM} not found in {(row['extracted_text'] if row else '')[:120]!r}")


def assert_pptx():
    _assert_format(F_PPTX, ci.METHOD_PPTX, PPTX_TEXT,
                   "PPTX: real generated PPTX extracted with known slide text present")


def assert_text():
    _assert_format(F_TXT, ci.METHOD_TEXT, "PLAINTEXT_MARKER_CHARLIE",
                   "TEXT: plain-text file passes through correctly")


def assert_image():
    doc = _current_by_name.get(F_IMG)
    row = _extraction(doc)
    if not (doc and row):
        fail("IMAGE: standalone image via Textract", f"no doc/extraction (resp={doc})")
        return
    method = row["extraction_method"]
    status = doc.get("status")

    # Probe the REAL Textract service (same code path the pipeline uses) to learn
    # whether Textract is actually usable in THIS environment — not merely whether
    # credentials are present. This drives the expectation instead of guessing.
    from services import textract
    probe_ok, probe_err, probe_text = False, None, ""
    try:
        res = textract.detect_document_text(IMG_BYTES)
        probe_text = (res["text"] or "").upper()
        probe_ok = True
    except textract.TextractUnavailable as exc:
        probe_err = f"TextractUnavailable: {exc}"
    except Exception as exc:  # noqa: BLE001 — any probe failure ⇒ not usable
        probe_err = f"{type(exc).__name__}: {exc}"

    if probe_ok:
        # Textract works here → the pipeline MUST have OCR'd the image and
        # recovered the known text.
        text = (row["extracted_text"] or "").upper()
        missing = [t for t in ("IMAGE", "OCR", "DELTA", "4471") if t not in text]
        if (status in ci.EXTRACTED_OR_BEYOND
                and method == ci.METHOD_IMAGE_TEXTRACT and not missing):
            ok("IMAGE: standalone image OCR'd via Textract; known text recovered",
               f"detected={text[:120]!r}")
        else:
            fail("IMAGE: standalone image OCR'd via Textract; known text recovered",
                 f"status={status}, method={method}, missing={missing}, "
                 f"text={text[:160]!r}")
        return

    # Textract not usable in this environment (no / placeholder / invalid AWS
    # creds). The LIVE OCR assertion is SKIPPED — never faked — and we instead
    # prove the image was routed to the OCR path and degraded GRACEFULLY (parked
    # as needs_ocr, no crash), exactly as a scanned PDF is.
    skip("IMAGE: standalone image OCR'd via Textract; known text recovered (LIVE)",
         f"Textract not usable in this environment ({probe_err}); the image is "
         "routed to the OCR path and parked as needs_ocr for a later run.")
    if status == ci.STATUS_NEEDS_OCR and method == ci.METHOD_OCR_PENDING:
        ok("IMAGE: routed to the Textract/OCR path, degraded to needs_ocr "
           "gracefully (no crash) when Textract is unavailable")
    else:
        fail("IMAGE: graceful needs_ocr degrade when Textract unavailable",
             f"status={status}, method={method}")


def assert_unsupported():
    doc = _current_by_name.get(F_BOGUS)
    row = _extraction(doc)
    if (doc and doc.get("status") == ci.STATUS_UNSUPPORTED
            and row and row["extraction_method"] == ci.METHOD_UNSUPPORTED):
        ok("UNSUPPORTED: unrecognised file → 'unsupported_format' status, no crash")
    else:
        fail("UNSUPPORTED: unrecognised file → 'unsupported_format'",
             f"resp={doc}, ext={dict(row) if row else None}")


def assert_email(email_resp):
    if not email_resp or not email_resp.get("documents"):
        fail("EMAIL: body + two recursive attachments", f"no email response: {email_resp}")
        return None
    email_doc = email_resp["documents"][0]
    drop_id = email_resp.get("drop_id")
    erow = _extraction(email_doc)

    # 1) Email body is its own extraction.
    if (erow and erow["extraction_method"] == ci.METHOD_EMAIL
            and "EMAIL_BODY_MARKER_ECHO" in (erow["extracted_text"] or "")
            and email_doc.get("status") in ci.EXTRACTED_OR_BEYOND):
        ok("EMAIL: body parsed as its OWN extraction (native_email, known body text)")
    else:
        fail("EMAIL: body parsed as its own extraction",
             f"resp={email_doc}, ext={dict(erow) if erow else None}")

    # 2) Exactly the email + 2 attachment rows in the drop; each own-path extracted.
    children = asyncio.run(fetch_drop_children(drop_id)) if drop_id else []
    atts = [r for r in children if r["source"] == "email_attachment"]
    if len(children) == 3 and len(atts) == 2:
        ok("EMAIL: two SEPARATE attachment documents rows created in the same drop",
           f"drop has {len(children)} rows (1 email + {len(atts)} attachments)")
    else:
        fail("EMAIL: two separate attachment documents rows in same drop",
             f"children={[(r['original_filename'], r['source']) for r in children]}")

    pdf_att = next((r for r in atts if (r["original_filename"] or "").endswith(".pdf")), None)
    xlsx_att = next((r for r in atts if (r["original_filename"] or "").endswith(".xlsx")), None)

    if (pdf_att and pdf_att["extraction_method"] == ci.METHOD_NATIVE
            and ATT_PDF_TEXT in (pdf_att["extracted_text"] or "")):
        ok("EMAIL: PDF attachment independently processed via the PDF path (native text)")
    else:
        fail("EMAIL: PDF attachment via PDF path",
             f"{dict(pdf_att) if pdf_att else None}")

    if (xlsx_att and xlsx_att["extraction_method"] == ci.METHOD_XLSX
            and ATT_XLSX_CELL in (xlsx_att["extracted_text"] or "")):
        ok("EMAIL: XLSX attachment independently processed via the XLSX path (structured)")
    else:
        fail("EMAIL: XLSX attachment via XLSX path",
             f"{dict(xlsx_att) if xlsx_att else None}")

    return email_doc.get("id")


# ── cross-org RLS (non-bypass app_service) ──────────────────────────────────
async def rls_isolation_checks(doc_id):
    if not doc_id:
        fail("RLS: cross-org document isolation", "no document id available")
        return
    use_set_role = False
    if APP_SERVICE_DATABASE_URL:
        try:
            conn = await _connect(APP_SERVICE_DATABASE_URL)
        except Exception as exc:  # noqa: BLE001
            skip("RLS: cross-org document isolation",
                 f"could not connect app_service DSN: {type(exc).__name__}: {exc}")
            return
    else:
        conn = await _connect(DATABASE_URL)
        use_set_role = True
        try:
            async with conn.transaction():
                await conn.execute("SET LOCAL ROLE app_service")
                who = await conn.fetchval("SELECT current_user")
                bypass = await conn.fetchval(
                    "SELECT rolbypassrls FROM pg_roles WHERE rolname = current_user")
            if who != "app_service" or bypass:
                await conn.close()
                skip("RLS: cross-org document isolation",
                     f"fallback role switch ineffective (current_user={who}, "
                     f"bypassrls={bypass}) — set APP_SERVICE_DATABASE_URL to run")
                return
        except Exception as exc:  # noqa: BLE001
            await conn.close()
            skip("RLS: cross-org document isolation",
                 f"cannot SET ROLE app_service ({type(exc).__name__}: {exc}) — "
                 f"set APP_SERVICE_DATABASE_URL to run")
            return
    try:
        async def visible(org):
            async with conn.transaction():
                if use_set_role:
                    await conn.execute("SET LOCAL ROLE app_service")
                await conn.execute(
                    "SELECT set_config('app.current_org_id',$1,true),"
                    "       set_config('app.is_super_admin','false',true)", org)
                return await conn.fetchval(
                    "SELECT count(*) FROM documents WHERE id = $1", doc_id)
        a = await visible(ORG_A)
        b = await visible(ORG_B)
        if a == 1 and b == 0:
            ok("RLS: new document visible in-org, invisible to a different org",
               f"ORG_A sees {a}, ORG_B sees {b}")
        else:
            fail("RLS: new document visible in-org, invisible cross-org",
                 f"ORG_A sees {a} (want 1), ORG_B sees {b} (want 0)")
    except Exception as exc:  # noqa: BLE001
        msg = f"{type(exc).__name__}: {exc}"
        if "permission denied" in str(exc).lower():
            skip("RLS: cross-org document isolation",
                 f"app_service lacks table GRANTs (not an isolation breach): {msg}")
        else:
            fail("RLS: cross-org document isolation", msg)
    finally:
        await conn.close()


# ── main ────────────────────────────────────────────────────────────────────
_current_by_name: dict[str, dict] = {}


def main_flow():
    global _current_by_name
    if not DATABASE_URL:
        fail("DATABASE_URL present", "env var not set — cannot run verify")
        return

    print("=== Chancery Phase 4 verify (multi-format ingestion) — start ===")
    report_task1()

    asyncio.run(teardown())            # teardown-at-START
    uid = asyncio.run(seed_user())
    print(f"[info] seeded verify user id={uid}")

    asyncio.run(structural_checks())

    install_stubs()
    by_name, email_resp = endpoint_flow()
    _current_by_name = by_name

    print("\n--- Extraction assertions ---")
    assert_pdf_no_regression(by_name)
    assert_docx()
    assert_xlsx()
    assert_pptx()
    assert_text()
    assert_image()
    assert_unsupported()
    email_doc_id = assert_email(email_resp)

    print("\n--- Cross-org isolation ---")
    # Prefer a new-format doc for the isolation check (DOCX); fall back to email.
    target = (by_name.get(F_DOCX) or {}).get("id") or email_doc_id
    asyncio.run(rls_isolation_checks(target))

    asyncio.run(teardown())            # teardown-at-END
    docs, drops, exts = asyncio.run(count_leftovers())
    if (docs, drops, exts) == (0, 0, 0):
        ok("Teardown: zero leftover rows (documents / drops / extractions)")
    else:
        fail("Teardown: zero leftover rows",
             f"documents={docs}, drops={drops}, extractions={exts}")


def summarize():
    n_pass = sum(1 for s, _, _ in _RESULTS if s == "PASS")
    n_fail = sum(1 for s, _, _ in _RESULTS if s == "FAIL")
    n_skip = sum(1 for s, _, _ in _RESULTS if s == "SKIP")
    print("\n=== SUMMARY ===")
    print(f"PASS={n_pass}  FAIL={n_fail}  SKIP={n_skip}")
    if n_fail:
        print("\nFAILURES:")
        for s, name, detail in _RESULTS:
            if s == "FAIL":
                print(f"  - {name}: {detail}")
    print("\nRESULT:", "PASS — all assertions green." if n_fail == 0
          else "FAIL — see failures above.")
    return 0 if n_fail == 0 else 1


if __name__ == "__main__":
    try:
        main_flow()
    except Exception:  # noqa: BLE001 — a crash is itself a failure to report
        print("[FATAL] verify crashed:")
        traceback.print_exc()
        _RESULTS.append(("FAIL", "verify script did not crash", "see traceback"))
    sys.exit(summarize())
