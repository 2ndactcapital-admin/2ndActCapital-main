"""Chancery intake — ROUTE + EXTRACT, now multi-format (Phase 4).

Phase 1 shipped a PDF-only pipeline: a deterministic native-text-layer ROUTE and
a pdfplumber native EXTRACT. Phase 2 added STORE (versioned R2) + SORT (open-set
classifier). Phase 4 GENERALIZES ROUTE into a real file-type dispatcher and adds
EXTRACT paths for the formats this platform actually receives — Word, Excel,
PowerPoint, plain text, standalone images, and email (body + attachments).

Design guarantees carried forward UNCHANGED for actual PDFs:
  * ROUTE is still a deterministic (no AI / no network) native-text-layer check.
  * A text-native PDF still goes to pdfplumber ``extract_native``; a scan/image
    PDF is still marked ``needs_ocr`` (a future OCR phase), never mis-extracted;
    a corrupt PDF still fails gracefully.

New in Phase 4:
  * ``detect_file_type`` inspects MAGIC BYTES / container structure — NOT the
    filename extension — so a mislabelled ``.pdf`` that is really an XLSX is
    routed as XLSX. Extension/MIME are only weak tie-breakers.
  * XLSX is recognised as already-structured data and routed straight to
    structured extraction (no "text layer" question applies).
  * DOCX / PPTX / plain-text / standalone-image extractors populate
    ``document_extractions`` the same way native PDF does (``extraction_method``
    reflecting the path used; ``extracted_text`` / ``extracted_tables`` as
    appropriate). Standalone images go through the shared AWS Textract service
    (``services.textract`` — the promoted, gate-proven DetectDocumentText call),
    NOT a duplicate integration.
  * Email intake parses the body as its own extraction AND creates a SEPARATE
    ``documents`` row per attachment (same ``document_drops`` batch), each run
    RECURSIVELY through this full pipeline — an attached PDF via the PDF path, an
    attached XLSX via the XLSX path, etc.
  * An unrecognised / unsupported file fails clearly (a ``documents`` row with
    status ``unsupported_format``) — never a crash, never a silent drop.

Bytes are handed in by the caller (the intake endpoint reads them); STORE
persists them to R2 when configured.
"""

import io
import json
import os
import re
import uuid as _uuid
import zipfile

import pdfplumber
from starlette.concurrency import run_in_threadpool

from services import storage, textract
from services.database import get_pool, set_rls_context, reset_rls_context
from services.document_classifier import classify_document

# extraction_method markers written to document_extractions.extraction_method
# (free text — no CHECK constraint on the column, confirmed via introspection).
METHOD_NATIVE = "native_pdfplumber"   # text-native PDF, real text extracted
METHOD_OCR_PENDING = "ocr_pending"    # scan / image-only → OCR (PDF scans, or
                                      # a standalone image when Textract is off)
METHOD_FAILED = "failed"              # corrupt / unreadable input
METHOD_PENDING = "pending"            # transient: routed, extraction not yet run
# Phase 4 method markers (one per new extraction path).
METHOD_DOCX = "native_docx"           # python-docx
METHOD_XLSX = "native_openpyxl"       # openpyxl (already-structured)
METHOD_PPTX = "native_pptx"           # python-pptx
METHOD_TEXT = "native_text"           # plain text, decoded
METHOD_IMAGE_TEXTRACT = "textract_image"  # standalone image OCR'd via Textract
METHOD_EMAIL = "native_email"         # email body (attachments get own rows)
METHOD_UNSUPPORTED = "unsupported"    # recognised as a file, but no handler

# documents.status values used by this phase (also free text, default 'dropped').
STATUS_ROUTED = "routed"
STATUS_EXTRACTED = "extracted"
STATUS_NEEDS_OCR = "needs_ocr"
STATUS_FAILED = "failed"
STATUS_UNSUPPORTED = "unsupported_format"  # Phase 4: unrecognised file type
# Phase 2 statuses (STORE + SORT). Extraction is NO LONGER terminal — STORE and
# SORT advance a successfully-extracted document further.
STATUS_STORED = "stored"                  # original bytes persisted to R2
STATUS_SORTED = "sorted"                  # classified into an existing category
STATUS_PENDING_REVIEW = "pending_review"  # classifier proposed a NEW category

# Any of these means "extraction succeeded" (the doc has real content).
EXTRACTED_OR_BEYOND = frozenset({
    STATUS_EXTRACTED, STATUS_STORED, STATUS_SORTED, STATUS_PENDING_REVIEW,
})

# Guard against pathological nesting (an email attached to an email attached to
# an email …). Depth 0 is the top-level dropped file.
_MAX_RECURSION_DEPTH = 6


# ---------------------------------------------------------------------------
# SORT — map a classifier doc_category code → Chancery doc_family
# ---------------------------------------------------------------------------
# The Chancery design splits documents into two downstream extraction families:
#   TABULAR    — value lives in tables / form fields (K-1s, tax returns,
#                financial statements, subscription & accreditation forms, IDs).
#   NARRATIVE  — value lives in prose legal instruments (formation docs, trusts,
#                wills, estate plans, operating agreements).
# There is NO reference_data 'doc_family' list in the deployed DB (confirmed by
# live introspection 2026-07-30: reference_data has the 12 canonical
# 'doc_category' rows and ZERO 'doc_family' rows), so this mapping is the
# code-level source of truth, keyed by those 12 codes. A code we do not
# recognise (e.g. a newly-ratified category) falls through to None rather than
# being force-fitted into a family.
FAMILY_TABULAR = "tabular"
FAMILY_NARRATIVE = "narrative"
FAMILY_OTHER = "other"

# The canonical doc_category code that triggers Phase-3 K-1 template extraction.
K1_CATEGORY = "k1"

_TABULAR_CATEGORIES = frozenset({
    K1_CATEGORY, "tax_return", "financial_statement", "subscription_doc",
    "accreditation", "id_document",
})
_NARRATIVE_CATEGORIES = frozenset({
    "llc_formation", "trust_instrument", "will", "estate_plan",
    "operating_agreement",
})


def doc_family_for_category(category_code: str | None) -> str | None:
    """Map a canonical doc_category code to its Chancery doc_family.

    Returns 'tabular' / 'narrative' for the two extraction families, 'other' for
    the explicit catch-all category, or None for an unknown/None code (so a code
    we cannot place is never guessed into a family).
    """
    if not category_code:
        return None
    if category_code in _TABULAR_CATEGORIES:
        return FAMILY_TABULAR
    if category_code in _NARRATIVE_CATEGORIES:
        return FAMILY_NARRATIVE
    if category_code == "other":
        return FAMILY_OTHER
    return None


# ---------------------------------------------------------------------------
# File-type detection — MAGIC BYTES / container structure, not the extension
# ---------------------------------------------------------------------------
# Canonical file-type tokens returned by ``detect_file_type`` and dispatched on
# by ``process_document``. Kept as constants so the router/verify never rely on
# a magic string.
TYPE_PDF = "pdf"
TYPE_DOCX = "docx"
TYPE_XLSX = "xlsx"
TYPE_PPTX = "pptx"
TYPE_TEXT = "text"
TYPE_IMAGE = "image"
TYPE_EML = "eml"
TYPE_MSG = "msg"
TYPE_UNSUPPORTED = "unsupported"

# OLE2 compound-file signature (legacy .msg / .doc / .xls all share it).
_OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
# ZIP local-file-header signatures (all OOXML — docx/xlsx/pptx — are ZIPs).
_ZIP_MAGICS = (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")


def _detect_ooxml(file_bytes: bytes) -> str:
    """Distinguish DOCX / XLSX / PPTX by looking INSIDE the ZIP container.

    OOXML files are ZIPs; the member paths are definitive (word/document.xml,
    xl/workbook.xml, ppt/presentation.xml) regardless of the filename. A ZIP that
    is none of these (e.g. a plain .zip) is unsupported.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            names = set(zf.namelist())
    except Exception:  # noqa: BLE001 — a broken ZIP is simply unsupported here
        return TYPE_UNSUPPORTED
    if "word/document.xml" in names:
        return TYPE_DOCX
    if "xl/workbook.xml" in names:
        return TYPE_XLSX
    if "ppt/presentation.xml" in names:
        return TYPE_PPTX
    # Fallback: some producers vary casing/paths — sniff [Content_Types].xml.
    if "[Content_Types].xml" in names:
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                ct = zf.read("[Content_Types].xml").decode("utf-8", "replace")
        except Exception:  # noqa: BLE001
            ct = ""
        if "wordprocessingml" in ct:
            return TYPE_DOCX
        if "spreadsheetml" in ct:
            return TYPE_XLSX
        if "presentationml" in ct:
            return TYPE_PPTX
    return TYPE_UNSUPPORTED


# RFC 5322 header lines we treat as evidence of an email. The "strong" set
# (rarely present verbatim in ordinary prose) gates the eml decision so a plain
# text file that merely contains a "Subject:" line is not misread as an email.
_EMAIL_STRONG_HEADERS = ("message-id:", "mime-version:", "received:",
                         "return-path:", "delivered-to:", "dkim-signature:")
_EMAIL_ADDR_HEADERS = ("from:", "to:", "cc:", "subject:", "date:", "sender:")


def _looks_like_email(text: str) -> bool:
    """Heuristic RFC-5322 detector over a decoded head of the file.

    Requires at least one 'strong' header (Message-ID / MIME-Version / Received /
    …) AND at least one addressing header (From / To / Subject / …) appearing as
    an actual header line near the top — enough to separate real ``.eml`` from
    prose that happens to mention "Subject:".
    """
    head_lines = text[:8192].splitlines()
    strong = addr = False
    for raw in head_lines[:60]:
        line = raw.lower()
        # A header line is "Name: value" with the name before a blank body line.
        if not raw.strip():
            # A blank line ends the header block; stop scanning at the body.
            if strong and addr:
                return True
            # keep scanning a little in case the very first line was blank
            continue
        if any(line.startswith(h) for h in _EMAIL_STRONG_HEADERS):
            strong = True
        if any(line.startswith(h) for h in _EMAIL_ADDR_HEADERS):
            addr = True
        if strong and addr:
            return True
    return strong and addr


def _is_mostly_printable(sample: bytes) -> bool:
    """True when a byte sample is overwhelmingly printable text (not binary)."""
    if not sample:
        return False
    if b"\x00" in sample:
        return False  # NUL byte ⇒ binary
    printable = sum(
        1 for b in sample
        if 32 <= b <= 126 or b in (9, 10, 13)  # printable + tab/newline/CR
    )
    return printable / len(sample) >= 0.85


def detect_file_type(
    file_bytes: bytes, filename: str | None = None, mime_type: str | None = None
) -> str:
    """Classify ``file_bytes`` by CONTENT (magic bytes / container), not name.

    Returns one of the ``TYPE_*`` tokens. The filename extension and MIME type
    are used only as weak tie-breakers (e.g. to prefer ``.msg`` reading on an OLE
    file) — a mislabelled extension can never override what the bytes actually
    are.
    """
    if not file_bytes:
        return TYPE_UNSUPPORTED
    head = file_bytes[:16]

    # --- Definitive binary signatures ---
    if head[:4] == b"%PDF" or head[:5] == b"%PDF-":
        return TYPE_PDF
    if head[:8] == _OLE_MAGIC:
        # Legacy OLE compound file. Among the formats this phase supports, only
        # Outlook .msg is OLE. The stream marker '__substg1.0_' is present in
        # .msg bodies (content-based), so we do not depend on the extension;
        # legacy .doc/.xls (also OLE) correctly fall through to unsupported.
        if b"__substg1.0_" in file_bytes or b"__properties_version1.0" in file_bytes:
            return TYPE_MSG
        if (filename or "").lower().endswith(".msg"):
            return TYPE_MSG
        return TYPE_UNSUPPORTED
    if head[:4] in _ZIP_MAGICS:
        return _detect_ooxml(file_bytes)

    # --- Image signatures ---
    if head[:8] == b"\x89PNG\r\n\x1a\n":
        return TYPE_IMAGE            # PNG
    if head[:3] == b"\xff\xd8\xff":
        return TYPE_IMAGE            # JPEG
    if head[:6] in (b"GIF87a", b"GIF89a"):
        return TYPE_IMAGE            # GIF
    if head[:4] in (b"II*\x00", b"MM\x00*"):
        return TYPE_IMAGE            # TIFF
    if head[:2] == b"BM":
        return TYPE_IMAGE            # BMP
    if head[:4] == b"RIFF" and file_bytes[8:12] == b"WEBP":
        return TYPE_IMAGE            # WEBP

    # --- Text-ish: email vs plain text ---
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = file_bytes.decode("latin-1")
        except Exception:  # noqa: BLE001
            return TYPE_UNSUPPORTED
    if _looks_like_email(text):
        return TYPE_EML
    if _is_mostly_printable(file_bytes[:4096]):
        return TYPE_TEXT
    return TYPE_UNSUPPORTED


# ---------------------------------------------------------------------------
# ROUTE — now a real dispatcher (PDF logic UNCHANGED for actual PDFs)
# ---------------------------------------------------------------------------
def route_document(
    file_bytes: bytes, *, filename: str | None = None, mime_type: str | None = None
) -> dict:
    """Route ``file_bytes`` by its actual type; returns a routing decision.

    Returns::

        {"file_type": str,          # a TYPE_* token
         "has_text_layer": bool | None,  # PDF-only; None for non-PDF
         "page_count": int,
         "valid": bool,             # False for unsupported/corrupt
         "valid_pdf": bool | None,  # PDF-only back-compat key
         "error": str | None}

    For an actual PDF this reproduces Phase-1 behaviour EXACTLY: at least one page
    with non-whitespace embedded text ⇒ ``has_text_layer=True``; a scan/image PDF
    ⇒ ``has_text_layer=False`` (route to OCR later); a corrupt/non-PDF that was
    sniffed as PDF ⇒ ``valid=False`` / ``valid_pdf=False`` with a clear error.
    XLSX (already-structured) reports ``has_text_layer=None`` — no text-layer
    question applies — and is routed straight to structured extraction downstream.
    """
    file_type = detect_file_type(file_bytes, filename, mime_type)
    result = {
        "file_type": file_type,
        "has_text_layer": None,
        "page_count": 0,
        "valid": True,
        "valid_pdf": None,
        "error": None,
    }

    if file_type == TYPE_PDF:
        # --- Phase-1 PDF ROUTE, verbatim behaviour ---
        try:
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                page_count = len(pdf.pages)
                has_text_layer = False
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    if text.strip():
                        has_text_layer = True
                        break
            result.update(
                has_text_layer=has_text_layer,
                page_count=page_count,
                valid=True,
                valid_pdf=True,
                error=None,
            )
        except Exception as exc:  # noqa: BLE001 — parse failure is a routing result, not a crash
            result.update(
                has_text_layer=False,
                page_count=0,
                valid=False,
                valid_pdf=False,
                error=f"{type(exc).__name__}: {exc}",
            )
        return result

    if file_type == TYPE_UNSUPPORTED:
        result["valid"] = False
        result["error"] = "unrecognised or unsupported file type"
        return result

    # Every other supported type is valid; the type-specific extractor decides
    # page_count / content. XLSX carries has_text_layer=None deliberately.
    return result


# ---------------------------------------------------------------------------
# EXTRACT (native PDF) — text + per-page tables for a text-native PDF
# ---------------------------------------------------------------------------
def extract_native(file_bytes: bytes) -> dict:
    """Extract full text and per-page tables from a TEXT-NATIVE PDF.

    Returns ``{"extracted_text": str, "extracted_tables": list, "page_count": int}``.
    ``extracted_tables`` is a list of ``{"page": int, "tables": [[...rows...]]}``
    entries, one per page that actually contained a table. Callers must only run
    this on a PDF that ``route_document`` reported as text-native — a scan will
    yield empty text here, which is why routing gates it.
    """
    page_texts: list[str] = []
    tables_by_page: list[dict] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        page_count = len(pdf.pages)
        for index, page in enumerate(pdf.pages, start=1):
            page_texts.append(page.extract_text() or "")
            page_tables = page.extract_tables() or []
            if page_tables:
                tables_by_page.append({"page": index, "tables": page_tables})
    return {
        "extracted_text": "\n\n".join(page_texts).strip(),
        "extracted_tables": tables_by_page,
        "page_count": page_count,
    }


# ---------------------------------------------------------------------------
# EXTRACT (Phase 4 formats) — each returns the same shape as extract_native
# ---------------------------------------------------------------------------
def _json_cell(value):
    """Coerce a spreadsheet/table cell to a JSON-serialisable, lossless form.

    Keep ints/floats/bools/None/str as-is; stringify anything else (dates,
    Decimals) so ``json.dumps`` never fails and numeric precision is preserved in
    text form (important for monetary figures — never silently reformatted)."""
    if value is None or isinstance(value, (int, float, bool, str)):
        return value
    return str(value)


def extract_docx(file_bytes: bytes) -> dict:
    """Extract paragraph text and tables from a DOCX via python-docx."""
    import docx  # local import — only needed on the DOCX path

    document = docx.Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    tables = []
    for t_index, table in enumerate(document.tables):
        rows = [[_json_cell(cell.text) for cell in row.cells] for row in table.rows]
        # Also fold table text into the body so classification/search see it.
        for row in rows:
            paragraphs.append(" | ".join(c for c in row if c))
        tables.append({"table": t_index, "rows": rows})
    return {
        "extracted_text": "\n".join(paragraphs).strip(),
        "extracted_tables": tables,
        "page_count": None,  # Word has no fixed page count without rendering
    }


def extract_xlsx(file_bytes: bytes) -> dict:
    """Extract every sheet of an XLSX as structured rows via openpyxl.

    XLSX is already-structured data, so ``extracted_tables`` is the primary
    output (one entry per sheet); ``extracted_text`` is a readable flattening of
    the same cells so the downstream text classifier still has something to read.
    ``data_only=True`` yields last-computed values rather than formulae.
    """
    import openpyxl  # local import — only needed on the XLSX path

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets = []
    text_parts = []
    try:
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                if row is None:
                    continue
                cells = [_json_cell(v) for v in row]
                # Trim wholly-empty trailing rows to keep output tight.
                if all(c is None or c == "" for c in cells):
                    continue
                rows.append(cells)
            sheets.append({"sheet": ws.title, "rows": rows})
            text_parts.append(f"# {ws.title}")
            for cells in rows:
                text_parts.append(
                    " | ".join("" if c is None else str(c) for c in cells)
                )
    finally:
        wb.close()
    return {
        "extracted_text": "\n".join(text_parts).strip(),
        "extracted_tables": sheets,
        "page_count": len(sheets),  # sheet count is the natural "unit" count
    }


def extract_pptx(file_bytes: bytes) -> dict:
    """Extract per-slide text and tables from a PPTX via python-pptx."""
    from pptx import Presentation  # local import — only on the PPTX path

    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts = []
    tables = []
    slide_count = 0
    for s_index, slide in enumerate(prs.slides, start=1):
        slide_count += 1
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if txt and txt.strip():
                    parts.append(txt.strip())
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = [
                    [_json_cell(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                tables.append({"slide": s_index, "rows": rows})
                for row in rows:
                    parts.append(" | ".join(c for c in row if c))
        if parts:
            slide_texts.append(f"[Slide {s_index}]\n" + "\n".join(parts))
    return {
        "extracted_text": "\n\n".join(slide_texts).strip(),
        "extracted_tables": tables,
        "page_count": slide_count,
    }


def extract_text(file_bytes: bytes) -> dict:
    """Decode a plain-text file (UTF-8, then latin-1 fallback)."""
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1", errors="replace")
    return {
        "extracted_text": text.strip(),
        "extracted_tables": [],
        "page_count": None,
    }


# ---------------------------------------------------------------------------
# Email parsing — body + attachments (attachments processed recursively later)
# ---------------------------------------------------------------------------
_TAG_RE = re.compile(r"<[^>]+>")
_WS_RE = re.compile(r"[ \t]*\n[ \t]*")


def _strip_html(html: str) -> str:
    """Very small HTML→text reduction for an HTML-only email body."""
    text = re.sub(r"(?is)<(script|style).*?</\1>", " ", html)
    text = _TAG_RE.sub(" ", text)
    text = (text.replace("&nbsp;", " ").replace("&amp;", "&")
            .replace("&lt;", "<").replace("&gt;", ">").replace("&quot;", '"'))
    text = _WS_RE.sub("\n", text)
    return re.sub(r"[ \t]{2,}", " ", text).strip()


def parse_email(file_bytes: bytes, file_type: str) -> dict:
    """Parse an email into its body text and a list of attachments.

    Returns ``{"body_text": str, "attachments": [{"filename", "content",
    "content_type"}]}``. ``.eml`` uses the Python stdlib ``email`` package (modern
    ``policy.default`` API); ``.msg`` uses the ``extract_msg`` library. Each
    attachment's raw bytes are returned for INDEPENDENT downstream processing —
    this function never extracts attachment content inline.
    """
    if file_type == TYPE_MSG:
        return _parse_msg(file_bytes)
    return _parse_eml(file_bytes)


def _parse_eml(file_bytes: bytes) -> dict:
    import email
    from email import policy

    msg = email.message_from_bytes(file_bytes, policy=policy.default)
    header_lines = []
    for key in ("Subject", "From", "To", "Cc", "Date"):
        value = msg.get(key)
        if value:
            header_lines.append(f"{key}: {value}")

    body_text = ""
    body = msg.get_body(preferencelist=("plain", "html"))
    if body is not None:
        try:
            content = body.get_content()
        except Exception:  # noqa: BLE001 — undecodable body ⇒ empty, not a crash
            content = ""
        if body.get_content_type() == "text/html":
            content = _strip_html(content)
        body_text = content

    attachments = []
    for part in msg.iter_attachments():
        try:
            payload = part.get_content()
        except Exception:  # noqa: BLE001
            payload = part.get_payload(decode=True) or b""
        if isinstance(payload, str):
            payload = payload.encode("utf-8", "replace")
        attachments.append({
            "filename": part.get_filename() or "attachment",
            "content": payload,
            "content_type": part.get_content_type(),
        })

    full_body = ("\n".join(header_lines) + "\n\n" + (body_text or "")).strip()
    return {"body_text": full_body, "attachments": attachments}


def _parse_msg(file_bytes: bytes) -> dict:
    import extract_msg

    msg = extract_msg.openMsg(io.BytesIO(file_bytes))
    try:
        header_lines = []
        for label, value in (("Subject", msg.subject), ("From", msg.sender),
                             ("To", msg.to), ("Date", getattr(msg, "date", None))):
            if value:
                header_lines.append(f"{label}: {value}")
        body_text = msg.body or ""

        attachments = []
        for att in msg.attachments:
            data = att.data
            if isinstance(data, str):
                data = data.encode("utf-8", "replace")
            if not isinstance(data, (bytes, bytearray)):
                # Nested .msg attachments expose a Message object, not bytes —
                # skip rather than fake; reported as a known gap in verify.
                continue
            fname = (att.getFilename() if hasattr(att, "getFilename") else None) \
                or getattr(att, "longFilename", None) \
                or getattr(att, "shortFilename", None) or "attachment"
            attachments.append({
                "filename": fname,
                "content": bytes(data),
                "content_type": None,
            })
    finally:
        msg.close()

    full_body = ("\n".join(header_lines) + "\n\n" + body_text).strip()
    return {"body_text": full_body, "attachments": attachments}


# ---------------------------------------------------------------------------
# STORE — persist original bytes to R2, versioned (reuses the S17 mechanism)
# ---------------------------------------------------------------------------
def _slugify(value: str) -> str:
    """Filesystem/URL-safe slug for a storage-key path segment."""
    slug = re.sub(r"[^A-Za-z0-9._-]+", "-", value or "").strip("-").lower()
    return slug[:80] or "file"


async def store_document(
    pool,
    document_id,
    org_id,
    file_bytes: bytes,
    *,
    original_filename: str | None,
    mime_type: str | None,
    entity_id=None,
) -> str | None:
    """Persist the ORIGINAL bytes to R2 (versioned) and record storage_key.

    Reuses the Sprint-17 R2 mechanism verbatim (``services.storage.upload_bytes``
    → the ``2ndactcapital-docs`` bucket) — NOT a second integration. The object
    key embeds ``document_id`` (unique per documents row), so a re-upload can
    NEVER overwrite a prior stored file. A human-meaningful, 1-based version
    number — counted over the (org, entity, filename) natural key across prior
    STORED rows — is embedded in the key path, mirroring ``entity_documents``'
    own version handling (each version is a distinct object; the prior is
    retained, never clobbered).

    Sets ``documents.storage_key`` and ``status = 'stored'``. Returns the
    storage_key, or None when R2 is not configured (unattended / CI) — in which
    case the document keeps its pre-store status rather than falsely claiming to
    be stored.
    """
    if not os.environ.get("R2_ACCOUNT_ID"):
        print(f"[chancery] store_document: R2 not configured — skipping STORE "
              f"for {document_id}")
        return None

    async with pool.acquire() as conn:
        version = await conn.fetchval(
            """
            SELECT count(*) + 1 FROM documents
            WHERE org_id = $1
              AND original_filename = $2
              AND coalesce(entity_id::text, '') = coalesce($3::text, '')
              AND storage_key IS NOT NULL
              AND id <> $4
            """,
            org_id, original_filename, entity_id, document_id,
        )

    stem, ext = os.path.splitext(original_filename or "")
    entity_seg = str(entity_id) if entity_id else "unfiled"
    storage_key = (
        f"chancery/{org_id}/{entity_seg}/{_slugify(stem)}/v{version}/"
        f"{document_id}{ext.lower()}"
    )

    # boto3 is synchronous — never block the event loop.
    await run_in_threadpool(storage.upload_bytes, storage_key, file_bytes, mime_type)

    async with pool.acquire() as conn:
        await conn.execute(
            "UPDATE documents SET storage_key = $1, status = $2, updated_at = now() "
            "WHERE id = $3",
            storage_key, STATUS_STORED, document_id,
        )
    return storage_key


# ---------------------------------------------------------------------------
# SORT — classify an extracted document, set doc_family / status
# ---------------------------------------------------------------------------
async def sort_document(pool, document_id, org_id, extracted_text: str) -> dict:
    """Classify an extracted document and set ``doc_family`` / ``status``.

    Calls the Sprint-25 open-set classifier (``document_classifier.
    classify_document``):
      * MATCH to an existing category → ``doc_family`` from
        ``doc_family_for_category``; ``status = 'sorted'``.
      * NEW-category proposal → classify_document has ALREADY queued the proposal
        into ``doc_category_proposals`` (the house 'AI proposes, human ratifies'
        review queue). We do NOT invent a doc_family: ``doc_family`` stays NULL
        and ``status = 'pending_review'``.
      * model unavailable (no API key) → classify returns a null category; we
        leave the document's status/doc_family untouched (still 'stored' /
        'extracted') rather than claiming a sort.

    Returns the classifier result dict, augmented with the resolved ``doc_family``.
    """
    async with pool.acquire() as conn:
        result = await classify_document(conn, org_id, extracted_text)

        category = result.get("category_code")
        is_new = result.get("is_new_proposal")

        if is_new:
            await conn.execute(
                "UPDATE documents SET status = $1, updated_at = now() WHERE id = $2",
                STATUS_PENDING_REVIEW, document_id,
            )
            result["doc_family"] = None
        elif category:
            family = doc_family_for_category(category)
            await conn.execute(
                "UPDATE documents SET doc_family = $1, status = $2, "
                "updated_at = now() WHERE id = $3",
                family, STATUS_SORTED, document_id,
            )
            result["doc_family"] = family
        else:
            # No category (model unavailable / failed) — do not claim a sort.
            result["doc_family"] = None
    return result


# ---------------------------------------------------------------------------
# Shared helpers used by every extraction path
# ---------------------------------------------------------------------------
async def _write_extraction(
    conn, extraction_id, *, method, extracted_text, extracted_tables, page_count
):
    """Populate the pre-created document_extractions row with real content."""
    await conn.execute(
        """
        UPDATE document_extractions
        SET extraction_method = $1,
            extracted_text = $2,
            extracted_tables = $3::jsonb,
            page_count = $4
        WHERE id = $5
        """,
        method,
        extracted_text,
        json.dumps(extracted_tables or []),
        page_count,
        extraction_id,
    )


async def _set_status(conn, document_id, status):
    await conn.execute(
        "UPDATE documents SET status = $1, updated_at = now() WHERE id = $2",
        status, document_id,
    )


async def _mark_failed(pool, document_id, extraction_id):
    async with pool.acquire() as conn:
        if extraction_id is not None:
            await conn.execute(
                "UPDATE document_extractions SET extraction_method = $1 WHERE id = $2",
                METHOD_FAILED, extraction_id,
            )
        await _set_status(conn, document_id, STATUS_FAILED)


async def _store_and_sort(pool, doc, org_id, file_bytes, extracted_text):
    """Run Phase-2 STORE then SORT, each degrading independently.

    STORE first so durable bytes are never lost, and so the TERMINAL status
    reflects the SORT outcome (a 'pending_review' proposal must not be clobbered
    by a later 'stored' write). A STORE or SORT failure logs and leaves the
    document at its last good status rather than losing the extraction.
    """
    document_id = doc["id"]
    try:
        await store_document(
            pool, document_id, org_id, file_bytes,
            original_filename=doc["original_filename"],
            mime_type=doc["mime_type"],
            entity_id=doc["entity_id"],
        )
    except Exception as exc:  # noqa: BLE001 — STORE failure ≠ pipeline failure
        print(f"[chancery] store_document failed for {document_id}: "
              f"{type(exc).__name__}: {exc}")

    sort_result = None
    try:
        sort_result = await sort_document(pool, document_id, org_id, extracted_text or "")
    except Exception as exc:  # noqa: BLE001 — SORT failure ≠ pipeline failure
        print(f"[chancery] sort_document failed for {document_id}: "
              f"{type(exc).__name__}: {exc}")

    # Phase 3: a K-1 (SORT → category='k1', doc_family='tabular') fires real
    # template extraction + Phase-5 auto-link as part of THIS same flow — not a
    # separate manual step.
    await _maybe_extract_k1(pool, doc, org_id, file_bytes, sort_result)


async def _maybe_extract_k1(pool, doc, org_id, file_bytes, sort_result):
    """When SORT classified the document as a K-1, run Phase-3 template
    extraction and hand off to Phase-5's REAL auto-link/propose logic.

    A K-1 extraction or linkage failure logs and leaves the document at its last
    good status rather than failing the whole pipeline (mirrors STORE/SORT)."""
    if not sort_result or sort_result.get("category_code") != K1_CATEGORY:
        return
    try:
        from services import textract_extraction  # local import: only on the K-1 path
        await textract_extraction.extract_k1_and_link(pool, doc, org_id, file_bytes)
    except Exception as exc:  # noqa: BLE001 — K-1 extraction failure ≠ pipeline failure
        print(f"[chancery] k1 extraction/link failed for {doc['id']}: "
              f"{type(exc).__name__}: {exc}")


# ---------------------------------------------------------------------------
# ORCHESTRATION — process a single already-persisted documents row
# ---------------------------------------------------------------------------
async def process_document(document_id, org_id, file_bytes: bytes, *, _depth: int = 0) -> None:
    """Route, then extract, a single ``documents`` row — any supported format.

    Steps:
      1. Load the ``documents`` row (validates it exists in ``org_id``).
      2. ``route_document`` (magic-byte dispatch) → write a
         ``document_extractions`` row recording the routing decision; set
         ``documents.status = 'routed'``.
      3. Dispatch on the detected type:
           * PDF        → Phase-1 native path (text-native extract, or 'needs_ocr'
                          for a scan, or 'failed' for a corrupt PDF) — UNCHANGED.
           * DOCX/XLSX/ → the matching Phase-4 native extractor; 'extracted'.
             PPTX/TEXT
           * IMAGE      → AWS Textract OCR ('extracted'); 'needs_ocr' when
                          Textract is unavailable.
           * EML/MSG    → parse body as this row's extraction, THEN create a
                          separate ``documents`` row per attachment (same drop)
                          and recurse through this function.
           * UNSUPPORTED→ 'unsupported_format' (clear, no crash, not dropped).
      4. On successful text extraction, run Phase-2 STORE + SORT.

    The org RLS context is set from ``org_id`` so this is safe to call standalone
    (a worker) and recursively (email attachments). Returns None.
    """
    tokens = set_rls_context(org_id, False)
    try:
        pool = await get_pool()

        # 1. Load the documents row (existence + in-org validation).
        async with pool.acquire() as conn:
            doc = await conn.fetchrow(
                "SELECT id, org_id, original_filename, mime_type, entity_id, "
                "status, drop_id, sequence_in_drop, created_by "
                "FROM documents WHERE id = $1 AND org_id = $2",
                document_id, org_id,
            )
            if doc is None:
                print(f"[chancery] process_document: document {document_id} "
                      f"not found in org {org_id}")
                return

        # 2. ROUTE (deterministic magic-byte dispatch — no AI, no network).
        routing = route_document(
            file_bytes,
            filename=doc["original_filename"],
            mime_type=doc["mime_type"],
        )
        file_type = routing["file_type"]

        # Record the routing decision on its OWN document_extractions row.
        async with pool.acquire() as conn:
            extraction_id = await conn.fetchval(
                """
                INSERT INTO document_extractions (
                    document_id, org_id, extraction_method,
                    has_native_text_layer, page_count
                ) VALUES ($1, $2, $3, $4, $5)
                RETURNING id
                """,
                document_id, org_id, METHOD_PENDING,
                routing["has_text_layer"],
                routing["page_count"] or None,
            )
            await _set_status(conn, document_id, STATUS_ROUTED)

        # 3. Dispatch on the routed file type.
        if file_type == TYPE_PDF:
            await _process_pdf(pool, doc, org_id, extraction_id, routing, file_bytes)
        elif file_type == TYPE_UNSUPPORTED:
            async with pool.acquire() as conn:
                await conn.execute(
                    "UPDATE document_extractions SET extraction_method = $1 WHERE id = $2",
                    METHOD_UNSUPPORTED, extraction_id,
                )
                await _set_status(conn, document_id, STATUS_UNSUPPORTED)
        elif file_type in (TYPE_EML, TYPE_MSG):
            await _process_email(
                pool, doc, org_id, extraction_id, file_bytes, file_type, _depth
            )
        elif file_type == TYPE_IMAGE:
            await _process_image(pool, doc, org_id, extraction_id, file_bytes)
        else:
            # DOCX / XLSX / PPTX / TEXT — synchronous native extractors.
            await _process_native_format(
                pool, doc, org_id, extraction_id, file_type, file_bytes
            )
    finally:
        reset_rls_context(tokens)


# ---------------------------------------------------------------------------
# Per-type handlers
# ---------------------------------------------------------------------------
async def _process_pdf(pool, doc, org_id, extraction_id, routing, file_bytes):
    """Phase-1 PDF branch — behaviour intentionally identical to Phase 1/2."""
    document_id = doc["id"]

    if not routing["valid_pdf"]:
        # Corrupt / unreadable — do NOT attempt extraction. Record failure.
        await _mark_failed(pool, document_id, extraction_id)
        return

    if routing["has_text_layer"]:
        try:
            native = extract_native(file_bytes)
        except Exception as exc:  # noqa: BLE001 — extraction failure ≠ batch failure
            print(f"[chancery] extract_native failed for {document_id}: "
                  f"{type(exc).__name__}: {exc}")
            await _mark_failed(pool, document_id, extraction_id)
            return

        async with pool.acquire() as conn:
            await _write_extraction(
                conn, extraction_id,
                method=METHOD_NATIVE,
                extracted_text=native["extracted_text"],
                extracted_tables=native["extracted_tables"],
                page_count=native["page_count"],
            )
            await _set_status(conn, document_id, STATUS_EXTRACTED)

        await _store_and_sort(pool, doc, org_id, file_bytes, native["extracted_text"])
        return

    # Valid PDF but no text layer → scan/image-only → future OCR phase.
    async with pool.acquire() as conn:
        await conn.execute(
            "UPDATE document_extractions SET extraction_method = $1 WHERE id = $2",
            METHOD_OCR_PENDING, extraction_id,
        )
        await _set_status(conn, document_id, STATUS_NEEDS_OCR)


_NATIVE_EXTRACTORS = {
    TYPE_DOCX: (extract_docx, METHOD_DOCX),
    TYPE_XLSX: (extract_xlsx, METHOD_XLSX),
    TYPE_PPTX: (extract_pptx, METHOD_PPTX),
    TYPE_TEXT: (extract_text, METHOD_TEXT),
}


async def _process_native_format(pool, doc, org_id, extraction_id, file_type, file_bytes):
    """DOCX / XLSX / PPTX / TEXT — run the extractor, record, STORE + SORT."""
    document_id = doc["id"]
    extractor, method = _NATIVE_EXTRACTORS[file_type]
    try:
        result = extractor(file_bytes)
    except Exception as exc:  # noqa: BLE001 — one file's failure must not stop the batch
        print(f"[chancery] {method} extraction failed for {document_id}: "
              f"{type(exc).__name__}: {exc}")
        await _mark_failed(pool, document_id, extraction_id)
        return

    async with pool.acquire() as conn:
        await _write_extraction(
            conn, extraction_id,
            method=method,
            extracted_text=result["extracted_text"],
            extracted_tables=result["extracted_tables"],
            page_count=result["page_count"],
        )
        await _set_status(conn, document_id, STATUS_EXTRACTED)

    await _store_and_sort(pool, doc, org_id, file_bytes, result["extracted_text"])


async def _process_image(pool, doc, org_id, extraction_id, file_bytes):
    """Standalone image → OCR via the shared AWS Textract service.

    When Textract is not configured/available, degrade to 'needs_ocr' (retryable)
    rather than failing — mirrors how a scanned PDF is parked for a later OCR run.
    """
    document_id = doc["id"]
    if not textract.textract_configured():
        async with pool.acquire() as conn:
            await conn.execute(
                "UPDATE document_extractions SET extraction_method = $1 WHERE id = $2",
                METHOD_OCR_PENDING, extraction_id,
            )
            await _set_status(conn, document_id, STATUS_NEEDS_OCR)
        return

    try:
        ocr = await run_in_threadpool(textract.detect_document_text, file_bytes)
    except textract.TextractUnavailable as exc:
        print(f"[chancery] textract unavailable for {document_id}: {exc}")
        async with pool.acquire() as conn:
            await conn.execute(
                "UPDATE document_extractions SET extraction_method = $1 WHERE id = $2",
                METHOD_OCR_PENDING, extraction_id,
            )
            await _set_status(conn, document_id, STATUS_NEEDS_OCR)
        return
    except Exception as exc:  # noqa: BLE001 — a genuine OCR fault ≠ batch failure
        print(f"[chancery] textract failed for {document_id}: "
              f"{type(exc).__name__}: {exc}")
        await _mark_failed(pool, document_id, extraction_id)
        return

    async with pool.acquire() as conn:
        await _write_extraction(
            conn, extraction_id,
            method=METHOD_IMAGE_TEXTRACT,
            extracted_text=ocr["text"],
            extracted_tables=[],
            page_count=1,
        )
        await _set_status(conn, document_id, STATUS_EXTRACTED)

    await _store_and_sort(pool, doc, org_id, file_bytes, ocr["text"])


async def _process_email(pool, doc, org_id, extraction_id, file_bytes, file_type, depth):
    """Email → body as this row's extraction, THEN a recursive doc per attachment."""
    document_id = doc["id"]
    try:
        parsed = parse_email(file_bytes, file_type)
    except Exception as exc:  # noqa: BLE001 — an unparseable email fails gracefully
        print(f"[chancery] parse_email failed for {document_id}: "
              f"{type(exc).__name__}: {exc}")
        await _mark_failed(pool, document_id, extraction_id)
        return

    body_text = parsed["body_text"]
    attachments = parsed["attachments"]

    # Record the email BODY as this document's own extraction.
    async with pool.acquire() as conn:
        await _write_extraction(
            conn, extraction_id,
            method=METHOD_EMAIL,
            extracted_text=body_text,
            extracted_tables=[],
            page_count=None,
        )
        await _set_status(conn, document_id, STATUS_EXTRACTED)

    # STORE the raw email bytes + SORT its body like any other document.
    await _store_and_sort(pool, doc, org_id, file_bytes, body_text)

    # Recurse into EACH attachment as an independently-processed document.
    if depth + 1 > _MAX_RECURSION_DEPTH:
        print(f"[chancery] max recursion depth reached for {document_id}; "
              f"{len(attachments)} attachment(s) not expanded")
        return

    for att in attachments:
        att_bytes = att.get("content") or b""
        if not att_bytes:
            continue
        child_id = _uuid.uuid4()
        async with pool.acquire() as conn:
            # Same drop as the parent email (that is the batch linkage); next
            # sequence position within the drop. If the parent has no drop
            # (standalone processing) the child simply has no drop either.
            next_seq = None
            if doc["drop_id"] is not None:
                next_seq = await conn.fetchval(
                    "SELECT coalesce(max(sequence_in_drop), 0) + 1 "
                    "FROM documents WHERE drop_id = $1",
                    doc["drop_id"],
                )
            await conn.execute(
                """
                INSERT INTO documents (
                    id, org_id, entity_id, original_filename, source, mime_type,
                    status, drop_id, sequence_in_drop, created_by
                ) VALUES ($1, $2, $3, $4, 'email_attachment', $5, 'dropped',
                          $6, $7, $8)
                """,
                child_id, org_id, doc["entity_id"],
                att.get("filename") or "attachment",
                att.get("content_type"), doc["drop_id"], next_seq,
                doc["created_by"],
            )

        # Full pipeline, recursively — the child's TYPE is re-detected from its
        # OWN bytes, so a PDF attachment takes the PDF path, an XLSX the XLSX
        # path, and a nested email recurses again (bounded by _MAX_RECURSION_DEPTH).
        try:
            await process_document(child_id, org_id, att_bytes, _depth=depth + 1)
        except Exception as exc:  # noqa: BLE001 — one attachment failing ≠ email failure
            print(f"[chancery] attachment processing crashed for child "
                  f"{child_id} of {document_id}: {type(exc).__name__}: {exc}")
            try:
                async with pool.acquire() as conn:
                    await _set_status(conn, child_id, STATUS_FAILED)
            except Exception as inner:  # noqa: BLE001
                print(f"[chancery] failed to mark child {child_id} failed: {inner}")
